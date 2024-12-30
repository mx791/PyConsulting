def _object_rels(obj):
    try:
        rels = obj.rels

        # Change required for python-pptx 0.6.22
        check_rels_content = [k for k in rels]
        if isinstance(check_rels_content.pop(), str):
            return [v for k, v in rels.items()]
        else:
            return [k for k in rels]
    except:
        return []


def _exp_add_slide(ppt, slide_layout):
    """
    Function to handle slide creation in the Presentation, to avoid issues caused by default implementation.

    :param slide_layout:
    :return:
    """

    def generate_slide_partname(self):
        """Return |PackURI| instance containing next available slide partname."""
        from pptx.opc.packuri import PackURI

        sldIdLst = self._element.get_or_add_sldIdLst()

        existing_rels = [k.target_partname for k in _object_rels(self)]
        partname_str = "/ppt/slides/slide%d.xml" % (len(sldIdLst) + 1)

        while partname_str in existing_rels:
            import random
            import string

            random_part = "".join(random.choice(string.ascii_letters) for i in range(2))
            partname_str = "/ppt/slides/slide%s%d.xml" % (
                random_part,
                len(sldIdLst) + 1,
            )

        return PackURI(partname_str)

    def add_slide_part(self, slide_layout):
        """
        Return an (rId, slide) pair of a newly created blank slide that
        inherits appearance from *slide_layout*.
        """
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT
        from pptx.parts.slide import SlidePart

        partname = generate_slide_partname(self)
        slide_layout_part = slide_layout.part
        slide_part = SlidePart.new(partname, self.package, slide_layout_part)
        rId = self.relate_to(slide_part, RT.SLIDE)
        return rId, slide_part.slide

    def add_slide_ppt(self, slide_layout):
        rId, slide = add_slide_part(self.part, slide_layout)
        slide.shapes.clone_layout_placeholders(slide_layout)
        self._sldIdLst.add_sldId(rId)
        return slide

    # slide_layout = self.get_master_slide_layout(slide_layout)
    return add_slide_ppt(ppt.slides, slide_layout)


def copy_shapes(source, dest):
    """
    Helper to copy shapes handling edge cases.

    :param source:
    :param dest:
    :return:
    """
    import copy

    from pptx.shapes.group import GroupShape

    # Copy all existing shapes
    for shape in source:
        if isinstance(shape, GroupShape):
            group = dest.shapes.add_group_shape()
            group.name = shape.name
            group.left = shape.left
            group.top = shape.top
            group.width = shape.width
            group.height = shape.height
            group.rotation = shape.rotation

            # Recursive copy of contents
            copy_shapes(shape.shapes, group)

            # Fix offset
            cur_el = group._element.xpath(".//p:grpSpPr")[0]
            ref_el = shape._element.xpath(".//p:grpSpPr")[0]
            parent = cur_el.getparent()
            parent.insert(parent.index(cur_el) + 1, copy.deepcopy(ref_el))
            parent.remove(cur_el)

            result = group
        elif hasattr(shape, "image"):
            import io

            # Get image contents
            content = io.BytesIO(shape.image.blob)
            result = dest.shapes.add_picture(
                content, shape.left, shape.top, shape.width, shape.height
            )
            result.name = shape.name
            result.crop_left = shape.crop_left
            result.crop_right = shape.crop_right
            result.crop_top = shape.crop_top
            result.crop_bottom = shape.crop_bottom
        elif hasattr(shape, "has_chart") and shape.has_chart:
            result = clone_chart(shape, dest)
        else:
            import copy

            newel = copy.deepcopy(shape.element)
            dest.shapes._spTree.insert_element_before(newel, "p:extLst")
            result = dest.shapes[-1]


def duplicate_slide(ppt, slide_index: int):
    """
    Duplicate the slide with the given number in presentation.
    Adds the new slide by default at the end of the presentation.

    :param ppt:
    :param slide_index: Slide number
    :return:
    """
    source = ppt.slides[slide_index]

    dest = _exp_add_slide(ppt, source.slide_layout)

    # Remove all shapes from the default layout
    for shape in dest.shapes:
        remove_shape(shape)

    # Copy all existing shapes
    copy_shapes(source.shapes, dest)

    # Copy existing references of known type
    # e.g. hyperlinks
    known_refs = [
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
    ]
    for rel in _object_rels(source.part):
        if rel.reltype in known_refs:
            if rel.is_external:
                dest.part.rels.get_or_add_ext_rel(rel.reltype, rel._target)
            else:
                dest.part.rels.get_or_add(rel.reltype, rel._target)

    # Copy all existing shapes
    if source.has_notes_slide:
        txt = source.notes_slide.notes_text_frame.text
        dest.notes_slide.notes_text_frame.text = txt

    return dest


def remove_shape(shape):
    """
    Helper to remove a specific shape.

    :source: https://stackoverflow.com/questions/64700638/is-there-a-way-to-delete-a-shape-with-python-pptx

    :param shape:
    :return:
    """
    el = shape.element  # --- get reference to XML element for shape
    el.getparent().remove(el)  # --- remove that shape element from its tree


### CHARTS

from typing import Union

import pandas as pd


def chart_to_dataframe(graphical_frame) -> pd.DataFrame:
    """
    Helper to parse chart data to a DataFrame.

    :source: https://openpyxl.readthedocs.io/en/stable/pandas.html

    :param graphical_frame:
    :return:
    """
    from io import BytesIO

    from openpyxl import load_workbook

    wb = load_workbook(
        BytesIO(graphical_frame.chart.part.chart_workbook.xlsx_part.blob),
        read_only=True,
    )

    ws = wb.active

    from itertools import islice

    import pandas as pd

    data = ws.values
    cols = next(data)[1:]
    data = list(data)
    idx = [r[0] for r in data]
    data = (islice(r, 1, None) for r in data)
    df = pd.DataFrame(data, index=idx, columns=cols)

    # Drop None columns
    return df.dropna(axis=1, how="all")


def dataframe_to_chart_data(df):
    """
    Transforms a DataFrame to a CategoryChartData for PPT compilation.

    The indexes of the DataFrame are the categories, with each column becoming a series.

    :param df:
    :return:
    """
    import numpy as np
    from pptx.chart.data import CategoryChartData

    copy_data = CategoryChartData()
    copy_data.categories = df.index.astype(str).to_list()

    edge_cases = 0
    for c in df.columns:
        series_data = df[c].copy()
        fixed_series_data = series_data.replace([np.inf, -np.inf, np.nan], None)

        edge_cases = edge_cases + np.count_nonzero(fixed_series_data != series_data)

        copy_data.add_series(str(c), fixed_series_data.to_list())

    # Warning over data filled for compatibility
    if edge_cases > 0:
        import warnings

        warnings.warn("Series data containing NaN/INF values: filled to empty")

    return copy_data


def update(shape, data):
    shape.chart.replace_data(dataframe_to_chart_data(data))

    # Fix for filling non category charts (XY, Bubble)
    id_attribute = (
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )

    chart_ref_id = shape.element.xpath(".//c:chart")[0].attrib[id_attribute]
    chart_part = shape.part.rels._rels[chart_ref_id].target_part
    chart_part._element.xpath(".//c:autoUpdate")[0].set("val", "1")

    point_series = chart_part._element.xpath(".//c:xVal")
    import copy

    for s in point_series:
        series_ref = s.getparent()

        # Find new reference (category information)
        x = series_ref.xpath(".//c:cat")[0]
        y = series_ref.xpath(".//c:val")[0]

        # Find existing reference (XY, Bubble)
        prev_x = series_ref.xpath(".//c:xVal")[0]
        prev_y = series_ref.xpath(".//c:yVal")[0]

        # Clean old contents
        [prev_x.remove(c) for c in prev_x.getchildren()]
        [prev_y.remove(c) for c in prev_y.getchildren()]

        # Add new contents
        [prev_x.append(c) for c in copy.deepcopy(x).getchildren()]
        [prev_y.append(c) for c in copy.deepcopy(y).getchildren()]

        # Remove category information
        series_ref.remove(x)
        series_ref.remove(y)


def clone_chart(graphical_frame, dest):
    """
    Helper to clone a chart with related styling.

    :param graphical_frame: General shape containing the .chart property
    :param dest: Shapes object on which to add the new chart
    :return:
    """
    chart = graphical_frame.chart

    df = chart_to_dataframe(graphical_frame)
    chart_data = dataframe_to_chart_data(df)

    new_chart = dest.shapes.add_chart(
        chart.chart_type,
        graphical_frame.left,
        graphical_frame.top,
        graphical_frame.width,
        graphical_frame.height,
        chart_data,
    )

    # Fix offset for Graphical shape
    import copy

    cur_el = new_chart._element.xpath(".//p:nvGraphicFramePr")[0]
    ref_el = graphical_frame._element.xpath(".//p:nvGraphicFramePr")[0]
    parent = cur_el.getparent()
    parent.insert(parent.index(cur_el) + 1, copy.deepcopy(ref_el))
    parent.remove(cur_el)

    # Clone styling from old chart to new one
    from random import randrange

    from lxml import etree
    from pptx.oxml import parse_xml

    id_attribute = (
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )

    old_chart_ref_id = graphical_frame.element.xpath(".//c:chart")[0].attrib[
        id_attribute
    ]
    chart_ref_id = new_chart.element.xpath(".//c:chart")[0].attrib[id_attribute]

    new_chart_part = new_chart.part.rels._rels[chart_ref_id].target_part
    old_chart_part = graphical_frame.part.rels._rels[old_chart_ref_id].target_part

    chart_data_reference_id = new_chart_part._element.xpath(".//c:externalData")[
        0
    ].attrib[id_attribute]

    cloned_styling = copy.deepcopy(old_chart_part._element)
    cloned_styling.xpath(".//c:externalData")[0].set(
        id_attribute, chart_data_reference_id
    )
    cloned_styling.xpath(".//c:autoUpdate")[0].set("val", "1")
    new_chart_part.part._element = cloned_styling

    # Parse other relationships of the chart
    from pptx.opc.constants import CONTENT_TYPE as CT
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from pptx.opc.package import XmlPart

    class ColorsPart(XmlPart):
        partname_template = "/ppt/charts/colors%d.xml"

        @classmethod
        def new(cls, package, element):
            part = cls.load(
                package.next_partname(cls.partname_template),
                CT.OFC_CHART_COLORS,
                package,
                element,
            )
            return part

    class StylePart(XmlPart):
        partname_template = "/ppt/charts/style%d.xml"

        @classmethod
        def new(cls, package, element):
            part = cls.load(
                package.next_partname(cls.partname_template),
                CT.OFC_CHART_STYLE,
                package,
                element,
            )
            return part

    new_chart_refs = new_chart_part.rels
    old_chart_refs = old_chart_part.rels

    # Fix styling and colors applied to the new chart
    for k, v in dict(old_chart_refs._rels).items():
        if (
            v.reltype
            == "http://schemas.microsoft.com/office/2011/relationships/chartStyle"
        ):
            targ = v.target_part

            new_el = parse_xml(copy.deepcopy(targ.blob))
            new_el.set("id", str(randrange(10**5, 10**9)))
            new_colors_ref = StylePart.new(targ.package, etree.tostring(new_el))
            new_chart_refs.get_or_add(
                "http://schemas.microsoft.com/office/2011/relationships/chartStyle",
                new_colors_ref,
            )
        elif v.reltype == RT.CHART_COLOR_STYLE:
            targ = v.target_part

            new_el = parse_xml(copy.deepcopy(targ.blob))
            new_el.set("id", str(randrange(10**5, 10**9)))
            new_colors_ref = ColorsPart.new(targ.package, etree.tostring(new_el))
            new_chart_refs.get_or_add(RT.CHART_COLOR_STYLE, new_colors_ref)

    return new_chart


### TABLES


def add_column(table):
    """
    Duplicates the last column of the table and appends it to the end.

    :param table: shape.table element
    """
    import copy

    from pptx.table import _Cell, _Column

    new_col = copy.deepcopy(table._tbl.tblGrid.gridCol_lst[-1])
    table._tbl.tblGrid.append(new_col)  # copies last grid element

    for tr in table._tbl.tr_lst:
        # duplicate last cell of each row
        new_tc = copy.deepcopy(tr.tc_lst[-1])

        # Fix for column styling
        last_tc = tr.xpath(".//a:tc")[-1]
        parent = last_tc.getparent()
        parent.insert(parent.index(last_tc) + 1, new_tc)

        # Clear new cell content
        cell = _Cell(new_tc, tr.tc_lst)
        cell.text_frame.clear()

    # Fix column not writable
    # https://stackoverflow.com/questions/64591452/using-copy-deepcopy-with-python-pptx-to-add-a-column-to-a-table-leads-to-cell-at
    from pptx import oxml

    for child in table._tbl.getchildren():
        if isinstance(child, oxml.table.CT_TableGrid):
            ws = set()
            for j in child:
                if j.w not in ws:
                    ws.add(j.w)
                else:
                    for elem in j:
                        j.remove(elem)

    # Create object in memory, in case some operations are done by the library
    col = _Column(new_col, table)


def remove_column(table, column_index: int):
    """
    Removes a specified column from the table.
    :param table: shape.table element

    """
    column = list(table.columns)[column_index]

    col_idx = table._tbl.tblGrid.index(column._gridCol)

    for tr in table._tbl.tr_lst:
        tr.remove(tr.tc_lst[col_idx])

    table._tbl.tblGrid.remove(column._gridCol)


def add_row(table) -> None:
    """
    Duplicates the last row and appends it to the end.
    :param table: shape.table element
    """
    import copy
    from random import randrange

    from pptx.table import _Cell, _Row

    new_row = copy.deepcopy(table._tbl.tr_lst[-1])  # copies last row element

    for tc in new_row.tc_lst:
        cell = _Cell(tc, new_row.tc_lst)
        cell.text = ""

    table._tbl.append(new_row)
    row = _Row(new_row, table)

    # Fix row not writable
    reference = row._tr.xpath(".//a:ext")[0]
    reference.getchildren()[0].set("val", str(randrange(10**5, 10**9)))


def remove_row(table, row_index: int) -> None:
    """
    Remove a specified row from the table.

    :param table: shape.table element
    :return:
    """
    row = list(table.rows)[row_index]

    table._tbl.remove(row._tr)


import copy
from random import randrange

from lxml import etree

### SLIDE MASTER & LAYOUT
from pptx.opc.constants import CONTENT_TYPE as CT
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import XmlPart
from pptx.opc.packuri import PackURI
from pptx.oxml import parse_xml
from pptx.parts.slide import SlideLayoutPart as SLP
from pptx.parts.slide import SlideMasterPart as SMP


class SlideLayoutPart(SLP):
    partname_template = "/ppt/slideLayouts/slideLayout%d.xml"

    @classmethod
    def new(cls, slide_masters, slide_master, element):
        existing = [
            i.target_partname
            for s in slide_masters
            for i in _object_rels(s.part)
            if "slideLayouts" in i.target_partname
        ]
        partname_str = cls.partname_template % (len(existing) + 1)

        while partname_str in existing:
            import random
            import string

            random_part = "".join(random.choice(string.ascii_letters) for i in range(2))
            partname_str = cls.partname_template % (random_part, len(existing) + 1)

        part = cls.load(
            PackURI(partname_str),
            CT.PML_SLIDE_LAYOUT,
            slide_master,
            element,
        )
        return part


class SlideMasterPart(SMP):
    partname_template = "/ppt/slideMasters/slideMaster%d.xml"

    @classmethod
    def new(cls, ppt, element):
        existing = [
            i.target_partname
            for i in _object_rels(ppt.part)
            if "slideMasters" in i.target_partname
        ]
        partname_str = cls.partname_template % (len(existing) + 1)

        while partname_str in existing:
            import random
            import string

            random_part = "".join(random.choice(string.ascii_letters) for i in range(2))
            partname_str = cls.partname_template % (random_part, len(existing) + 1)

        part = cls.load(
            PackURI(partname_str),
            CT.PML_SLIDE_MASTER,
            ppt,
            element,
        )
        return part


class ThemePart(XmlPart):
    partname_template = "/ppt/theme/theme%d.xml"

    @classmethod
    def new(cls, ppt, element):
        existing = [
            i.target_partname
            for i in _object_rels(ppt.part)
            if "theme" in i.target_partname
        ]
        partname_str = cls.partname_template % (len(existing) + 1)

        while partname_str in existing:
            import random
            import string

            random_part = "".join(random.choice(string.ascii_letters) for i in range(2))
            partname_str = cls.partname_template % (random_part, len(existing) + 1)

        part = cls.load(
            PackURI(partname_str),
            CT.OFC_THEME,
            ppt,
            element,
        )
        return part


def _clone_sml_shapes(source, dest):
    # Remove all shapes from the default layout
    for shape in dest.shapes:
        remove_shape(shape)

    # Copy all existing shapes
    copy_shapes(source.shapes, dest)

    # Copy existing references of known type
    # e.g. hyperlinks
    known_refs = [
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
    ]
    for rel in _object_rels(source.part):
        if rel.reltype in known_refs:
            if rel.is_external:
                dest.part.rels.get_or_add_ext_rel(rel.reltype, rel._target)
            else:
                dest.part.rels.get_or_add(rel.reltype, rel._target)


def _new_existing_slide_ml_id(ppt):
    existing_slide_layout_ids = []
    for sm in ppt.slide_masters:
        existing_slide_layout_ids = existing_slide_layout_ids + sm.element.xpath(
            ".//p:sldLayoutId/@id"
        )

    existing_slide_layout_ids = existing_slide_layout_ids + ppt.element.xpath(
        ".//p:sldMasterId/@id"
    )
    sel_id = max([255] + [int(id_str) for id_str in existing_slide_layout_ids]) + 1

    return sel_id


def clone_slide_master(pres, slide_master):
    # Generate the new Slide Master part object
    new_el = copy.deepcopy(slide_master.element)
    new_ref = SlideMasterPart.new(pres.part, etree.tostring(new_el))

    # Connect the presentation to the new Slide Master
    rId = pres.part.relate_to(new_ref, RT.SLIDE_MASTER)

    # Add the Slide Master to the list of the available ones in the presentation
    sel_id = _new_existing_slide_ml_id(pres)
    el_ref = pres.slide_masters._sldMasterIdLst._add_sldMasterId()
    el_ref.set("id", str(sel_id))
    el_ref.set(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id", rId
    )

    # Remove references of Slide Layouts from the new Slide Master
    master_slide_ref = pres.slide_masters[-1]
    master_slide_ref = _fix_package_ref(master_slide_ref)  # Fixes for shape checks
    for i in master_slide_ref.element.xpath(".//p:sldLayoutId"):
        i.getparent().remove(i)

    # Copy known references relevant for the Slide Master
    known_refs = [
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
    ]
    for rel in _object_rels(slide_master.part):
        if rel.reltype in known_refs:
            if rel.is_external:
                master_slide_ref.part.rels.get_or_add_ext_rel(rel.reltype, rel._target)
            else:
                master_slide_ref.part.rels.get_or_add(rel.reltype, rel._target)

        # Themes need to be copied completely
        if rel.reltype == RT.THEME:
            targ = rel.target_part

            new_el = parse_xml(copy.deepcopy(targ.blob))
            new_el.set("id", str(randrange(10**5, 10**9)))

            new_ref = ThemePart.new(pres, etree.tostring(new_el))
            master_slide_ref.part.rels.get_or_add(RT.THEME, new_ref)
            pres.part.rels.get_or_add(RT.THEME, new_ref)

    # Ensure all shapes of the Slide Master are set up correctly
    _clone_sml_shapes(slide_master, master_slide_ref)

    return master_slide_ref


def clone_slide_layout(ppt, source_layout, dest_slide_master):
    """
    Clones an existing slide layout to a destination slide master.
    Expected to be used with the clone_slide_master utility.
    """
    slide_master = dest_slide_master

    # Generate the new Slide Layout part object
    new_el = copy.deepcopy(source_layout.element)
    new_ref = SlideLayoutPart.new(
        ppt.slide_masters, slide_master.part, etree.tostring(new_el)
    )

    # Connect the Slide Master to the new Slide Layout
    rId = slide_master.part.relate_to(new_ref, RT.SLIDE_LAYOUT)

    # Add the Slide Master to the list of the available ones in the Slide Master
    sel_id = _new_existing_slide_ml_id(ppt)
    el_ref = slide_master.slide_layouts._sldLayoutIdLst._add_sldLayoutId()
    el_ref.set("id", str(sel_id))
    el_ref.set(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id", rId
    )

    # Connect the Slide Layout to the Slide Master
    dest = slide_master.slide_layouts[-1]
    dest = _fix_package_ref(dest)  # Fixes for shape checks

    dest.part.relate_to(slide_master.part, RT.SLIDE_MASTER)

    # Ensure all shapes of the Slide Layout are set up correctly
    _clone_sml_shapes(source_layout, dest)

    return None


from pptx.package import Package
from pptx.shapes.shapetree import _BaseGroupShapes
from pptx.slide import SlideLayout, SlideMaster


class _BaseGroupShapesProxy(_BaseGroupShapes):
    def __init__(self, el):
        super().__init__(el._spTree, el)


def _fix_package_ref(dest):
    class SlideMasterFix(SlideMaster):
        def __init__(self, el):
            self._shapes = None
            super().__init__(el._element, el._part)

        @property
        def shapes(self):
            if self._shapes is None:
                self._shapes = _BaseGroupShapesProxy(super().shapes)

            return self._shapes

    class SlideLayoutFix(SlideLayout):
        def __init__(self, el):
            self._shapes = None
            super().__init__(el._element, el._part)

        @property
        def shapes(self):
            if self._shapes is None:
                self._shapes = _BaseGroupShapesProxy(super().shapes)

            return self._shapes

    if isinstance(dest, SlideMaster):
        dest = SlideMasterFix(dest)
    else:
        dest = SlideLayoutFix(dest)

    # Fix for Slide Layout not having a direct connection to the PPT package
    pack = dest.part.package
    if not isinstance(pack, Package):
        pack = pack.package

    if not isinstance(pack, Package):
        pack = pack.package

    dest.part._package = pack

    return dest


### EXPERIMENTS ON TEXT SIZE


def estimate_text_box_size(
    txt, font, max_width: Union[int, None] = None, line_spacing: int = 4  # ImageFont
):
    """
    Example of use:
    right_margin = left_margin = Length(Cm(0.25)).pt * pt_per_pixel
    top_margin = bottom_margin = Length(Cm(0.13)).pt * pt_per_pixel
    width, height = estimate_text_box_size(
        txt,
        font,
        max_width=shape_width - (right_margin + left_margin),
    )

    print("Computed in pixels (w, h)")
    print((width + right_margin + left_margin, height + top_margin + bottom_margin))


    :param txt:
    :param font:
    :param max_width:
    :param line_spacing:
    :return:
    """

    from PIL import Image, ImageDraw

    image = Image.new(size=(400, 300), mode="RGB")
    draw = ImageDraw.Draw(image)
    emu_per_inch = 914400
    px_per_inch = 72.0
    pt_per_pixel = 0.75

    fontsize_pt = 12
    # font = ImageFont.truetype("arial.ttf", int(fontsize_pt / pt_per_pixel))
    import math
    import textwrap

    if max_width is not None:
        actual_txt = []
        for line in txt.split("\n"):
            _, _, width, h = font.getbbox(line)
            split_at = len(line) // math.ceil(width / max_width)
            actual_txt = actual_txt + textwrap.wrap(line, width=split_at)

        new_lines = len(actual_txt)
        actual_txt = "\n".join(actual_txt)
    else:
        actual_txt = txt
        new_lines = 0

    left, top, right, bottom = draw.multiline_textbbox(
        (0, 0), actual_txt, font=font, spacing=line_spacing
    )
    ascent, descent = font.getmetrics()

    return right - left, bottom  # + descent * new_lines