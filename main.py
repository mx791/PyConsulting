from pptx import Presentation 
import datetime
import yfinance as yf
from utils import duplicate_slide
import pandas as pd
from pptx.chart.data import ChartData
import re

def handle_title_slide(pres: Presentation):
    date = datetime.datetime.now()
    pres.slides[1].shapes[0].text = date.strftime('%d/%m/%Y')


def update_summary(pres: Presentation, indexes: list[int]):
    pres.slides[1].shapes[1].text = "\n".join(indexes)



def add_company_description(pres: Presentation, name: str, summary: str, desription: str):
    new_slide = duplicate_slide(pres, 2)
    new_slide.shapes[0].text = name

    new_description = desription
    for i in re.findall("([.]) ([A-Z])", desription):
        replace = i[0] + " " + i[1]
        new_description = new_description.replace(replace, ".\n" + i[1])
    
    #new_description = [new_description.replace(" ".join(m), "\n") for m in ]
    new_slide.shapes[1].text_frame.paragraphs[0].runs[0].text = new_description
    new_slide.shapes[2].text_frame.paragraphs[0].runs[0].text = summary


def add_kpi(pres: Presentation, name: str, kpis: dict):
    new_slide = duplicate_slide(pres, 3)
    new_slide.shapes[0].text = name
    for id, name in enumerate(kpis):
        if id >= 6:
            break
        new_slide.shapes[1+id*2].text_frame.paragraphs[0].runs[0].text = kpis[name]
        new_slide.shapes[2+id*2].text_frame.paragraphs[0].runs[0].text = name


def format_number(number: float, abbreviate: bool = True) -> str:
    if number > 10**9:
        return f"{int(number / 10**9)} B" if abbreviate else f"{(number / 10**9):2f} BB"
    
    if number > 10**6:
        return f"{int(number / 10**6)} M" if abbreviate else f"{(number / 10**6):2f} M"
    
    if number > 10**3:
        return f"{int(number / 10**3)} K" if abbreviate else f"{(number / 10**3):2f} K"

    return f"{int(number)}" if abbreviate else f"{(number):2f}"


def add_stock_price(pres: Presentation, name: str, data: pd.DataFrame):
    new_slide = duplicate_slide(pres, 5)
    new_slide.shapes[1].text = name

    chart_data = ChartData()
    chart_data.categories = pd.to_datetime(data.index.values)
    chart_data.add_series('Stock', tuple(data["Close"].values))

    chart = new_slide.shapes[0].chart
    chart.replace_data(chart_data)

    #chart.has_legend = True
    #chart.legend.include_in_layout = False
    chart.series[0].smooth = True


def add_financials(pres: Presentation, name, data):
    new_slide = duplicate_slide(pres, 4)

    new_slide.shapes[3].text = name

    values = tuple(reversed([data[c].values[-1] / 1_000_000_000 for c in data.columns[:-1]]))
    chart_data = ChartData()
    chart_data.categories = tuple(reversed([c.strftime('%d/%m/%Y') for c in data.columns[:-1]]))
    chart_data.add_series('Stock', values)

    chart = new_slide.shapes[0].chart
    chart.replace_data(chart_data)


    values = tuple(reversed([data[c].values[-4] / 1_000_000_000 for c in data.columns[:-1]]))
    chart_data = ChartData()
    chart_data.categories = tuple(reversed([c.strftime('%d/%m/%Y') for c in data.columns[:-1]]))
    chart_data.add_series('Stock', values)

    chart = new_slide.shapes[1].chart
    chart.replace_data(chart_data)





def company_data(pres: Presentation, symbol: str):
    ticker  = yf.Ticker(symbol)
    add_company_description(
        pres,
        ticker.info["shortName"],
        "Secteur: " + ticker.info["sector"],
        ticker.info["longBusinessSummary"]
    )
    add_kpi(
        pres,
        ticker.info["shortName"],
        {
            "Forward PE": f'{ticker.info["forwardPE"]:.2f}',
            "Trailing PE": f'{ticker.info["trailingPE"]:.2f}',
            "Market Value": format_number(ticker.info["marketCap"]) + "$",
            "Share price": f'{ticker.info["currentPrice"]:.2f}$',
            "Earnings": format_number(ticker.info["freeCashflow"]) + "$",
            "Gross Margin": f'{ticker.info["grossMargins"]:.2f}',
        }
    )
    add_financials(pres, ticker.info["shortName"], ticker.financials)
    add_stock_price(pres, ticker.info["shortName"], ticker.history("3Y"))


if __name__ == "__main__":

    ppt = Presentation("./SP500 Updates.pptx")
    handle_title_slide(ppt)

    symbols = [
        "AAPL", "MSFT", "NVDA", "TSLA", "AMZN", "AVGO"
    ]

    update_summary(ppt, symbols)
    for s in symbols:
        company_data(ppt, s)


    ppt.slides[2]

    ppt.save("./out.pptx")