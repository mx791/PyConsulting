from pptx import Presentation


def process_shortname(name: str) -> str:
    """Text processing for company names."""
    return name.split("  ")[0]


def replace_in_slide(slide, key_values: dict[str, str]):
    """Template matching function for all texts in a slide."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                for key in key_values:
                    if f"[{key}]" in r.text:
                        r.text = r.text.replace(f"[{key}]", key_values[key])
                    elif f"[{key}]" in p.text:
                        p.text = p.text.replace(f"[{key}]", key_values[key])


def format_number(number: float, abbreviate: bool = True) -> str:
    """Shorter a large number to make it more human readable."""
    minus = number < 0
    number = number if not minus else -number
    if number > 10**9:
        return f"{int(number / 10**9)} B" if abbreviate else f"{(number / 10**9):2f} B"

    if number > 10**6:
        return f"{int(number / 10**6)} M" if abbreviate else f"{(number / 10**6):2f} M"

    if number > 10**3:
        return f"{int(number / 10**3)} K" if abbreviate else f"{(number / 10**3):2f} K"

    return f"{"-" if minus else ""}{int(number)}" if abbreviate else f"{(number):2f}"


def delete_slide(presentation: Presentation, index: int):
    """Delete a slide given its id in the presentation."""
    xml_slides = presentation.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[index])
