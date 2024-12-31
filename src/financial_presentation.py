import datetime

import yfinance as yf
import pandas as pd
from pptx import Presentation

from utils_ext import duplicate_slide
from pptx.chart.data import ChartData
from llm import create_company_summary, top_companies_to_invest
from utils import replace_in_slide, format_number, process_shortname, delete_slide


def add_title_slide(pres: Presentation, speaker: str):
    """Creates the title slide."""
    slide = duplicate_slide(pres, 0)
    date = datetime.datetime.now()
    replace_in_slide(
        slide,
        {
            "date": date.strftime("%d/%m/%Y"),
            "speaker": speaker,
        },
    )


def add_summary_slide(pres: Presentation, indexes: list[str]):
    """Add summary slide, with company names."""
    slide = duplicate_slide(pres, 1)
    new_content = "\n".join([f"- {v}" for v in indexes])
    replace_in_slide(slide, {"table_of_content": new_content})


def add_company_intro_slide(pres: Presentation, name: str, subtitle: str):
    """Add company name slide."""
    slide = duplicate_slide(pres, 2)
    replace_in_slide(slide, {"company_name": name, "subtitle": subtitle})


def add_company_description_slide(pres: Presentation, name: str):
    """Add company description slide."""
    new_slide = duplicate_slide(pres, 3)
    description = create_company_summary(name)
    replace_in_slide(new_slide, {"company_name": name, "content": description})


def add_kpi_slide(pres: Presentation, name: str, kpis: dict):
    """Add KPI slide."""
    new_slide = duplicate_slide(pres, 4)
    new_slide.shapes[0].text = name

    target_dict = {"company_name": name}

    for id, name in enumerate(kpis):
        if id >= 6:
            break
        target_dict[f"kpi_{id+1}"] = name
        target_dict[f"kpi_{id+1}_value"] = kpis[name]

    replace_in_slide(new_slide, target_dict)


def get_kpis_data(ticker: yf.Ticker):
    """Extract KPIs data."""
    return {
        "Forward PE": f'{ticker.info["forwardPE"]:.2f}',
        "Trailing PE": f'{ticker.info["trailingPE"]:.2f}',
        "Market Value": format_number(ticker.info["marketCap"]) + "$",
        "Share price": f'{ticker.info["currentPrice"]:.2f}$',
        "Earnings": format_number(ticker.info["freeCashflow"]) + "$",
        "Gross Margin": f'{ticker.info["grossMargins"]:.2f}',
    }


def add_financials_results_slide(pres: Presentation, name: str, data):
    """Creates the revenues slide."""
    new_slide = duplicate_slide(pres, 5)

    replace_in_slide(new_slide, {"company_name": name})
    one_bn = 1_000_000_000

    values = tuple(reversed([data[c].values[-1] / one_bn for c in data.columns[:-1]]))
    chart_data = ChartData()
    chart_data.categories = tuple(
        reversed([c.strftime("%d/%m/%Y") for c in data.columns[:-1]])
    )
    chart_data.add_series("Stock", values)

    chart = new_slide.shapes[2].chart
    chart.replace_data(chart_data)

    values = tuple(reversed([data[c].values[-4] / one_bn for c in data.columns[:-1]]))
    chart_data = ChartData()
    chart_data.categories = tuple(
        reversed([c.strftime("%d/%m/%Y") for c in data.columns[:-1]])
    )
    chart_data.add_series("Stock", values)

    chart = new_slide.shapes[3].chart
    chart.replace_data(chart_data)


def add_stock_price_slide(pres: Presentation, name: str, data: pd.DataFrame):
    """Creates the stock price slide."""
    new_slide = duplicate_slide(pres, 6)
    replace_in_slide(new_slide, {"company_name": name})

    chart_data = ChartData()
    chart_data.categories = pd.to_datetime(data.index.values)
    chart_data.add_series("Stock", tuple(data["Close"].values))

    chart = new_slide.shapes[2].chart
    chart.replace_data(chart_data)

    chart.series[0].smooth = True


def add_recommendation_slide(pres, companies):
    """Creates the recommendation slide."""
    new_slide = duplicate_slide(pres, 7)
    content = top_companies_to_invest(companies)
    content = content.replace("\n\n", "\n").split("\n")
    target_dict = {}
    for i in range(len(content)):
        target_dict[f"company_{i+1}"] = content[i].split("] :")[0][1:]
        target_dict[f"explain_{i+1}"] = content[i].split("] :")[1]
    replace_in_slide(new_slide, target_dict)


def make_financial_pres(companies_symbols, speaker_info, output_path):
    """Creates the presentation."""

    # fetching company data from yahoo finance
    ticker_list, company_names = [], []
    for symbol in companies_symbols:
        try:
            ticker = yf.Ticker(symbol)
            ticker.history()
            ticker_list.append(ticker)
            company_names.append(process_shortname(ticker.info["shortName"]))
            print(f"Fetched data for {symbol} : {company_names[-1]}")
        except Exception:
            print(
                f"Company symbol: {symbol} does not appear in the yahoo finance database"
            )

    # opening template file
    ppt = Presentation("./template/template.pptx")
    nb_of_slide = len(ppt.slides)

    # update first slides
    add_title_slide(ppt, speaker_info)
    add_summary_slide(ppt, company_names)

    for ticker, name in zip(ticker_list, company_names):
        add_company_intro_slide(ppt, name, f"Sector: {ticker.info["sector"]}")
        add_company_description_slide(ppt, name)
        add_kpi_slide(ppt, name, get_kpis_data(ticker))
        add_financials_results_slide(ppt, name, ticker.financials)
        add_stock_price_slide(ppt, name, ticker.history("2y"))

    add_recommendation_slide(ppt, company_names)

    for _ in range(nb_of_slide):
        delete_slide(ppt, 0)

    ppt.save(output_path)
    print(f"Final presentation available at {output_path}")
